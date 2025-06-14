# -*- coding: utf-8 -*-

import os
import sys
import tempfile
import shutil
import logging
from pathlib import Path
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Cm
from pptx.enum.shapes import MSO_SHAPE_TYPE
from matplotlib.ticker import FuncFormatter
from openpyxl import load_workbook
import re
from .variaveis import VARIAVEIS_PLANILHA, GRAFICOS, ESTILO_GRAFICOS, POSICOES_GRAFICOS, VARIAVEIS_SLIDES
from copy import deepcopy

# Configuração de logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

class Config:
    """Gerencia configurações globais"""
    def __init__(self):
        self.base_dir = Path(__file__).parent.parent.parent
        
        # Configurar diretórios
        self.input_dir = self.base_dir / 'input'
        self.output_dir = self.base_dir / 'output'
        self.templates_dir = self.base_dir / 'templates'
        
        # Criar diretórios necessários
        for dir_path in [self.input_dir, self.output_dir, self.templates_dir]:
            dir_path.mkdir(exist_ok=True)
        
        # Carregar configurações dos gráficos do arquivo de mapeamento
        self.graph = ESTILO_GRAFICOS

class ResourceManager:
    """Gerencia recursos temporários"""
    def __init__(self):
        self.temp_dir = None
        self.resources = []
    
    def __enter__(self):
        self.temp_dir = Path(tempfile.mkdtemp(prefix="solar_proposta_"))
        logger.info(f"Criado diretório temporário: {self.temp_dir}")
        return self
    
    def __exit__(self, exc_type, exc_val, exc_tb):
        if self.temp_dir and self.temp_dir.exists():
            shutil.rmtree(self.temp_dir)
            logger.info(f"Removido diretório temporário: {self.temp_dir}")
    
    def get_path(self, filename):
        """Retorna caminho para um arquivo temporário"""
        return self.temp_dir / filename

class ExcelReader:
    """Lê e processa dados do Excel"""
    def __init__(self, excel_path, config):
        self.path = Path(excel_path)
        self.config = config
        self.workbook = None
    
    def __enter__(self):
        try:
            self.workbook = load_workbook(self.path, data_only=True)
            return self
        except Exception as e:
            logger.error(f"Erro ao abrir Excel: {e}")
            raise
    
    def __exit__(self, exc_type, exc_val, exc_tb):
        if self.workbook:
            self.workbook.close()
    
    def read_variables(self):
        """Lê variáveis da aba 'Variáveis'"""
        try:
            sheet = self.workbook["Variáveis"]
            variables = {}
            
            # Log do conteúdo da aba
            logger.info(f"Conteúdo da aba 'Variáveis':")
            for row in sheet.iter_rows(min_row=1, max_row=2, min_col=1, max_col=sheet.max_column):
                for cell in row:
                    logger.info(f"  {cell.coordinate}: {cell.value}")
            
            logger.info("=== Iniciando leitura de variáveis ===")
            logger.info(f"Dimensões da aba 'Variáveis': {sheet.dimensions}")
            
            # Ler variáveis usando o mapeamento
            for var_name, cell_coord in VARIAVEIS_PLANILHA.items():
                cell = sheet[cell_coord]
                if cell.value is not None:
                    # Se a célula estiver além da linha 2, pular
                    if cell.row > 2:
                        continue
                    variables[var_name] = str(cell.value)
                    logger.info(f"✓ Variável encontrada: {var_name} = {cell.value} (célula {cell_coord})")
                else:
                    logger.warning(f"✗ Variável {var_name} sem valor definido (célula {cell_coord})")
            
            logger.info(f"Total de variáveis encontradas: {len(variables)}")
            logger.info("=== Fim da leitura de variáveis ===\n")
            return variables
        except Exception as e:
            logger.error(f"Erro ao ler variáveis: {e}")
            raise
    
    def read_graph_data(self):
        """Lê dados para gráficos da aba 'Tabelas Graficos Proposta'"""
        try:
            sheet = self.workbook["Tabelas Graficos Proposta"]
            data = {}
            
            logger.info("=== Iniciando leitura de dados para gráficos ===")
            logger.info(f"Dimensões da aba 'Tabelas Graficos Proposta': {sheet.dimensions}")
            
            # Ler dados de cada gráfico usando o mapeamento
            for graph_num, graph_config in GRAFICOS.items():
                logger.info(f"\nLendo dados para {graph_config['titulo']}:")
                graph_data = {}
                
                try:
                    for key, cell_range in graph_config['ranges'].items():
                        # Extrair valores das células
                        values = []
                        for row in sheet[cell_range]:
                            value = row[0].value
                            if value is not None:
                                try:
                                    # Tentar converter para float se for número
                                    if isinstance(value, (int, float)):
                                        values.append(float(value))
                                    else:
                                        # Se for string, manter como está
                                        values.append(str(value))
                                except (ValueError, TypeError):
                                    values.append(str(value))
                            else:
                                # Usar 0 para valores vazios em dados numéricos
                                if key not in ['mes', 'ano']:
                                    values.append(0.0)
                                else:
                                    values.append('')
                        
                        graph_data[key] = values
                        non_none_values = [v for v in values if v is not None and v != '']
                        logger.info(f"- {key} ({cell_range}): {len(non_none_values)} valores não nulos de {len(values)} total")
                    
                    data[graph_num] = graph_data
                    
                except Exception as e:
                    logger.error(f"Erro ao ler dados para {graph_config['titulo']}: {e}")
                    # Criar dados vazios para este gráfico
                    data[graph_num] = {
                        key: [0.0] * (12 if key == 'mes' else 25)
                        for key in graph_config['ranges'].keys()
                    }
            
            logger.info("=== Fim da leitura de dados para gráficos ===\n")
            return data
            
        except Exception as e:
            logger.error(f"Erro ao ler dados dos gráficos: {e}")
            # Retornar estrutura vazia mas válida
            return {
                graph_num: {
                    key: [0.0] * (12 if key == 'mes' else 25)
                    for key in graph_config['ranges'].keys()
                }
                for graph_num, graph_config in GRAFICOS.items()
            }

class GraphManager:
    """Gerencia criação e formatação de gráficos"""
    def __init__(self, config, resource_manager):
        self.config = config
        self.resources = resource_manager
        self.setup_style()
        
        # Mapeamento de cores
        self.cores = {
            'vermelho': self.config.graph['cores']['vermelho'],
            'verde': self.config.graph['cores']['verde']
        }
    
    def get_color(self, nome_cor):
        """Retorna o código hexadecimal da cor"""
        return self.cores.get(nome_cor.lower(), '#000000')  # Preto como cor padrão
    
    def setup_style(self):
        """Configura estilo padrão dos gráficos"""
        plt.style.use('default')
        plt.rcParams.update({
            'figure.figsize': [
                self.config.graph['dimensoes']['largura'] * self.config.graph['dimensoes']['cm_para_polegadas'],
                self.config.graph['dimensoes']['altura'] * self.config.graph['dimensoes']['cm_para_polegadas']
            ],
            'figure.dpi': self.config.graph['geral']['dpi'],
            'savefig.dpi': self.config.graph['geral']['dpi'],
            'font.size': 10,  # Tamanho da fonte padrão
            'axes.grid': self.config.graph['geral']['mostrar_grade'],
            'axes.spines.top': self.config.graph['geral']['mostrar_borda']['topo'],
            'axes.spines.right': self.config.graph['geral']['mostrar_borda']['direita'],
            'figure.constrained_layout.use': True,
            'figure.autolayout': False,
            'figure.facecolor': 'none',  # Fundo transparente
            'axes.facecolor': 'none',    # Fundo transparente
            'savefig.facecolor': 'none', # Fundo transparente ao salvar
            'savefig.transparent': True  # Garante transparência ao salvar
        })
    
    def format_currency(self, x, p=None):
        """Formata valores monetários para o matplotlib"""
        if abs(x) >= 100000:
            # Formato 1k para valores acima de 100 mil
            return f'R${x/1000:.0f}k'
        else:
            # Formato R$1000 para valores menores
            return f'R${x:,.0f}'.replace(',', '.')
    
    def add_data_labels(self, x, y, color, offset=0):
        """Adiciona labels aos pontos do gráfico"""
        # Lista de índices para adicionar labels (primeiro, a cada 5 anos, e último)
        indices = [0] + list(range(4, len(x), 5))  # Primeiro e a cada 5 anos
        if len(x) - 1 not in indices:  # Adiciona o último se ainda não estiver incluído
            indices.append(len(x) - 1)
        
        for i in indices:
            plt.annotate(
                self.format_currency(y[i]),
                (x[i], y[i]),
                xytext=(0, offset),
                textcoords='offset points',
                ha='center',
                va='bottom',
                fontsize=12,
                color='black'
            )
    
    def save_graph(self, name):
        """Salva gráfico como imagem"""
        path = self.resources.get_path(f"{name}.png")
        plt.savefig(path, bbox_inches='tight', pad_inches=0.1)
        plt.close()
        return path
    
    def create_graphs(self, data):
        """Cria gráficos com os dados fornecidos"""
        graph_paths = {}
        
        try:
            for graph_id, graph_config in GRAFICOS.items():
                plt.figure()
                
                # Configurar grid horizontal para gráficos específicos
                if graph_id in ['graph1', 'graph3', 'graph4', 'graph5']:
                    plt.grid(True, axis='y', linestyle='--', alpha=0.5)
                
                # Configurar grid vertical para gráfico 2
                if graph_id == 'graph2':
                    plt.grid(True, axis='x', linestyle='--', alpha=0.5)
                
                if graph_config['tipo'] == 'linha':
                    # Gráfico de linha simples
                    plt.plot(data[graph_id]['ano'], 
                           data[graph_id]['valor'],
                           color=self.get_color(graph_config['cores']['valor']))
                    # Mostrar todos os anos
                    plt.xticks(data[graph_id]['ano'])
                    # Adicionar labels a cada 5 anos
                    if graph_id == 'graph2':
                        self.add_data_labels(
                            data[graph_id]['ano'],
                            data[graph_id]['valor'],
                            self.get_color(graph_config['cores']['valor'])
                        )
                
                elif graph_config['tipo'] == 'linhas':
                    # Gráfico de múltiplas linhas
                    plt.plot(data[graph_id]['ano'],
                           data[graph_id]['economia'],
                           label='Economia',
                           color=self.get_color(graph_config['cores']['economia']))
                    plt.plot(data[graph_id]['ano'],
                           data[graph_id]['custo'],
                           label='Custo',
                           color=self.get_color(graph_config['cores']['custo']))
                    plt.legend(fontsize=10)
                    # Mostrar todos os anos
                    plt.xticks(data[graph_id]['ano'])
                    # Adicionar labels a cada 5 anos
                    if graph_id == 'graph5':
                        self.add_data_labels(
                            data[graph_id]['ano'],
                            data[graph_id]['economia'],
                            self.get_color(graph_config['cores']['economia']),
                            offset=10
                        )
                        self.add_data_labels(
                            data[graph_id]['ano'],
                            data[graph_id]['custo'],
                            self.get_color(graph_config['cores']['custo']),
                            offset=-10
                        )
                
                elif graph_config['tipo'] == 'barras':
                    if 'valor' in data[graph_id]:
                        # Gráfico de barras simples
                        plt.bar(data[graph_id]['ano'], 
                               data[graph_id]['valor'],
                               color=self.get_color(graph_config['cores']['valor']))
                        # Mostrar todos os anos
                        plt.xticks(data[graph_id]['ano'])
                    elif 'positivo' in data[graph_id]:
                        # Gráfico de barras duplas (positivo/negativo)
                        plt.bar(data[graph_id]['ano'], 
                               data[graph_id]['positivo'],
                               color=self.get_color(graph_config['cores']['positivo']))
                        plt.bar(data[graph_id]['ano'], 
                               data[graph_id]['negativo'],
                               color=self.get_color(graph_config['cores']['negativo']))
                        # Mostrar todos os anos
                        plt.xticks(data[graph_id]['ano'])
                
                elif graph_config['tipo'] == 'barras_duplas':
                    # Gráfico de barras duplas (produção vs consumo)
                    x = range(len(data[graph_id]['mes']))
                    width = 0.35
                    
                    plt.bar([i - width/2 for i in x], 
                           data[graph_id]['producao'],
                           width,
                           label='Produção',
                           color=self.get_color(graph_config['cores']['producao']))
                    
                    plt.bar([i + width/2 for i in x], 
                           data[graph_id]['consumo'],
                           width,
                           label='Consumo',
                           color=self.get_color(graph_config['cores']['consumo']))
                    
                    plt.xticks(x, data[graph_id]['mes'])
                    plt.legend(fontsize=10)
                
                # Configurar título e labels
                plt.title(graph_config['titulo'], fontsize=10)
                plt.xlabel(graph_config['eixos']['x'], fontsize=10)
                plt.ylabel(graph_config['eixos']['y'], fontsize=10)
                
                # Configurar formatação do eixo Y para valores monetários
                if 'Valor (R$)' in graph_config['eixos']['y']:
                    plt.gca().yaxis.set_major_formatter(FuncFormatter(self.format_currency))
                
                # Configurar tamanho dos ticks
                plt.xticks(fontsize=10)
                plt.yticks(fontsize=10)
                
                # Rotacionar labels do eixo x para melhor legibilidade
                plt.xticks(rotation=45)
                
                # Ajustar layout para evitar corte de labels
                plt.tight_layout()
                
                # Salvar gráfico
                graph_paths[graph_id] = self.save_graph(f"graph_{graph_id}")
                
        except Exception as e:
            logger.error(f"Erro ao criar gráficos: {e}")
            raise
        
        return graph_paths

class PresentationManager:
    """Gerencia criação e edição da apresentação"""
    def __init__(self, template_path, config):
        self.presentation = Presentation(template_path)
        self.config = config
    
    def insert_graph(self, slide_index, graph_path, graph_id):
        """Insere gráfico em um slide"""
        try:
            slide = self.presentation.slides[slide_index]
            
            # Remover placeholder existente
            for shape in slide.shapes:
                if shape.shape_type in [MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.PLACEHOLDER]:
                    if hasattr(shape, 'placeholder_format'):
                        if shape.placeholder_format.type in [13, 18]:
                            shape.element.getparent().remove(shape.element)
            
            # Pegar posição do gráfico
            pos = POSICOES_GRAFICOS[graph_id]
            
            # Adicionar gráfico na posição específica
            picture = slide.shapes.add_picture(
                str(graph_path),
                Cm(pos['left']),
                Cm(pos['top']),
                Cm(pos['width']),
                Cm(pos['height'])
            )
            
        except Exception as e:
            logger.error(f"Erro ao inserir gráfico no slide {slide_index + 1}: {e}")
            raise
    
    def find_variables_in_shape(self, shape, slide_idx):
        """Encontra todas as variáveis em um shape"""
        pattern = re.compile(r'\{\{(\w+)\}\}')
        found_vars = set()
        
        if hasattr(shape, "text_frame"):
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    matches = pattern.finditer(run.text)
                    for match in matches:
                        var_name = match.group(1)
                        found_vars.add((var_name, run.text))
        
        if hasattr(shape, "table"):
            for row_idx, row in enumerate(shape.table.rows):
                for col_idx, cell in enumerate(row.cells):
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            matches = pattern.finditer(run.text)
                            for match in matches:
                                var_name = match.group(1)
                                found_vars.add((var_name, run.text))
        
        return found_vars
    
    def replace_variables(self, variables):
        """Substitui variáveis no texto mantendo a formatação original"""
        try:
            logger.info("=== Iniciando análise do template PowerPoint ===")
            all_vars = set()
            
            # Criar um dicionário com as variáveis em minúsculas e maiúsculas
            variables_normalized = {}
            for var_name, value in variables.items():
                # Preservar quebras de linha
                if isinstance(value, str):
                    value = value.replace('\\n', '\n')
                variables_normalized[var_name.lower()] = value
                variables_normalized[var_name.upper()] = value
                variables_normalized[var_name] = value  # Manter o caso original também
            
            for slide_idx, slide in enumerate(self.presentation.slides):
                logger.info(f"\nAnalisando Slide {slide_idx + 1}:")
                
                # Obter lista de variáveis esperadas para este slide
                expected_vars = set(VARIAVEIS_SLIDES.get(slide_idx, []))
                if expected_vars:
                    logger.info(f"Variáveis esperadas no Slide {slide_idx + 1}: {', '.join(expected_vars)}")
                
                # Encontrar todas as variáveis no slide
                variables_to_replace = []
                found_vars = set()  # Conjunto para rastrear variáveis encontradas
                
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame"):
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                # Procurar por padrões de variável no texto
                                matches = re.finditer(r'{{([^}]+)}}', run.text)
                                for match in matches:
                                    var_name = match.group(1)
                                    logger.info(f"Encontrada variável no texto: {var_name} em {run.text}")
                                    variables_to_replace.append((run, var_name, match.group(0)))
                                    found_vars.add(var_name.lower())
                
                # Verificar se todas as variáveis esperadas foram encontradas
                if expected_vars:
                    expected_vars_lower = {v.lower() for v in expected_vars}
                    missing_vars = expected_vars_lower - found_vars
                    if missing_vars:
                        logger.warning(f"\n⚠️ Variáveis não encontradas no Slide {slide_idx + 1}:")
                        for var in missing_vars:
                            logger.warning(f"  - {var}")
                    
                    extra_vars = found_vars - expected_vars_lower
                    if extra_vars:
                        logger.warning(f"\n⚠️ Variáveis encontradas mas não esperadas no Slide {slide_idx + 1}:")
                        for var in extra_vars:
                            logger.warning(f"  - {var}")
                
                # Substituir as variáveis encontradas
                substituted_vars = set()
                for run, var_name, original_text in variables_to_replace:
                    # Tentar substituir com o caso exato, depois minúsculo, depois maiúsculo
                    replaced = False
                    for test_name in [var_name, var_name.lower(), var_name.upper()]:
                        if test_name in variables:
                            value = str(variables[test_name])
                            logger.info(f"Tentando substituir {test_name} com valor: {value}")
                            
                            # Encontrar o shape que contém este run
                            for shape in slide.shapes:
                                if hasattr(shape, "text_frame"):
                                    # Verificar se o texto do shape contém a variável
                                    if original_text in shape.text:
                                        # Se o texto contém quebras de linha
                                        if "\n" in value:
                                            # Dividir o texto em linhas
                                            lines = value.split("\n")
                                            # Limpar o texto frame
                                            shape.text_frame.clear()
                                            # Adicionar cada linha como um novo parágrafo
                                            for line in lines:
                                                p = shape.text_frame.add_paragraph()
                                                p.text = line
                                                # Copiar a formatação do run original
                                                for new_run in p.runs:
                                                    new_run.font.size = run.font.size
                                                    new_run.font.bold = run.font.bold
                                                    new_run.font.italic = run.font.italic
                                                    new_run.font.underline = run.font.underline
                                                    if hasattr(run.font, 'color') and run.font.color:
                                                        new_run.font.color.rgb = run.font.color.rgb
                                        else:
                                            # Se não há quebras de linha, substituir mantendo a formatação
                                            for paragraph in shape.text_frame.paragraphs:
                                                if original_text in paragraph.text:
                                                    paragraph.text = value
                                                    # Copiar a formatação do run original
                                                    for new_run in paragraph.runs:
                                                        new_run.font.size = run.font.size
                                                        new_run.font.bold = run.font.bold
                                                        new_run.font.italic = run.font.italic
                                                        new_run.font.underline = run.font.underline
                                                        if hasattr(run.font, 'color') and run.font.color:
                                                            new_run.font.color.rgb = run.font.color.rgb
                                        
                                        replaced = True
                                        substituted_vars.add(var_name.lower())
                                        logger.info(f"✓ Substituída variável {var_name} = {value}")
                                        break
                            if replaced:
                                break
                    
                    if not replaced:
                        logger.warning(f"✗ Variável {var_name} não encontrada para substituição")
                    
                    # Registrar para o log final
                    all_vars.add((var_name, original_text))
                
                # Verificar se todas as variáveis esperadas foram substituídas
                if expected_vars:
                    missing_substitutions = expected_vars_lower - {v.lower() for v in substituted_vars}
                    if missing_substitutions:
                        logger.error(f"\n❌ Variáveis não substituídas no Slide {slide_idx + 1}:")
                        for var in missing_substitutions:
                            logger.error(f"  - {var}")
                        raise ValueError(f"Falha ao substituir todas as variáveis no Slide {slide_idx + 1}")
                
                # Log do slide
                if variables_to_replace:
                    logger.info(f"\nResumo do Slide {slide_idx + 1}:")
                    for run, var_name, original_text in variables_to_replace:
                        var_exists = any(var_name.lower() == v.lower() for v in expected_vars)
                        status = "✓" if var_exists and var_name.lower() in substituted_vars else "✗"
                        logger.info(f"{status} {{{{{var_name}}}}} (Original: {original_text.strip()})")
            
            # Log final
            logger.info("\nResumo de todas as variáveis encontradas no template:")
            for var_name, original_text in sorted(all_vars):
                var_exists = any(var_name.lower() == k.lower() for k in variables.keys())
                status = "✓" if var_exists else "✗"
                value = next((v for k, v in variables.items() if k.lower() == var_name.lower()), "NÃO ENCONTRADA")
                logger.info(f"{status} {{{{{var_name}}}}} = {value}")
            
            logger.info("=== Fim da substituição de variáveis ===\n")
        except Exception as e:
            logger.error(f"Erro ao substituir variáveis: {e}")
            raise
    
    def _copy_run_format(self, source_run, target_run):
        """Copia a formatação de um run para outro"""
        target_run.font.name = source_run.font.name
        target_run.font.size = source_run.font.size
        target_run.font.bold = source_run.font.bold
        target_run.font.italic = source_run.font.italic
        target_run.font.underline = source_run.font.underline
        if hasattr(source_run.font, 'color') and source_run.font.color:
            target_run.font.color.rgb = source_run.font.color.rgb
    
    def save(self, output_path):
        """Salva apresentação"""
        self.presentation.save(output_path)

def main(excel_path=None, template_path=None, output_path=None, verbose=False):
    """Função principal"""
    try:
        # Configurar logging
        if verbose:
            logging.getLogger().setLevel(logging.INFO)
        else:
            logging.getLogger().setLevel(logging.WARNING)
        
        # Verificar argumentos da linha de comando se não fornecidos
        if excel_path is None or template_path is None or output_path is None:
            if len(sys.argv) != 4:
                print("Uso: python gerar_apresentacao.py <excel> <template> <saida>")
                return 1
            excel_path = Path(sys.argv[1])
            template_path = Path(sys.argv[2])
            output_path = Path(sys.argv[3])
        
        # Converter para Path se necessário
        excel_path = Path(excel_path)
        template_path = Path(template_path)
        output_path = Path(output_path)
        
        # Verificar arquivos
        if not excel_path.exists():
            raise FileNotFoundError(f"Excel não encontrado: {excel_path}")
        if not template_path.exists():
            raise FileNotFoundError(f"Template não encontrado: {template_path}")
        
        # Inicializar configuração
        config = Config()
        
        # Processar apresentação
        with ResourceManager() as resources:
            # Ler dados do Excel
            with ExcelReader(excel_path, config) as excel:
                variables = excel.read_variables()
                graph_data = excel.read_graph_data()
            
            # Gerar gráficos
            graph_mgr = GraphManager(config, resources)
            graph_paths = graph_mgr.create_graphs(graph_data)
            
            # Criar apresentação
            prs_mgr = PresentationManager(template_path, config)
            
            # Inserir gráficos usando as posições definidas em POSICOES_GRAFICOS
            for graph_id, pos in POSICOES_GRAFICOS.items():
                prs_mgr.insert_graph(pos['slide'], graph_paths[graph_id], graph_id)
            
            # Substituir variáveis
            prs_mgr.replace_variables(variables)
            
            # Salvar resultado
            prs_mgr.save(output_path)
        
        print("✅ Apresentação gerada com sucesso!")
        return 0
    
    except Exception as e:
        print(f"❌ Erro: {str(e)}")
        return 1

if __name__ == "__main__":
    sys.exit(main())
    
